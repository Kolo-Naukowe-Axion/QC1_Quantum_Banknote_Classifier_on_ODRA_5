#!/usr/bin/env python3
"""Restructure slides 16-19 to present QPU and Simulator KL methodologies in parallel."""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "Dania - IEEE QCCL 2026 (draft) 2.pptx.pptx"
BACKUP = ROOT / "Dania - IEEE QCCL 2026 (draft) 2.pptx.backup.pptx"

TITLE_COLOR = RGBColor(0x19, 0x1E, 0x3A)
BODY_COLOR = RGBColor(0x00, 0x00, 0x00)
ACCENT_QPU = RGBColor(0x21, 0x25, 0x31)
ACCENT_SIM = RGBColor(0xC4, 0x30, 0x2B)
MUTED_COLOR = RGBColor(0x44, 0x44, 0x44)


def find_shape(slide, name: str):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def set_single_line_text(shape, text: str, *, bold: bool | None = None, size_pt: float = 26) -> None:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.color.rgb = BODY_COLOR
    if bold is not None:
        run.font.bold = bold


def set_title(shape, text: str) -> None:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(42)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR


def set_footer(shape, text: str) -> None:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(13)
    run.font.color.rgb = RGBColor(0xE4, 0xE4, 0xE4)


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    lines: list[str],
    *,
    size_pt: float = 22,
    bold_first: bool = False,
    color: RGBColor = BODY_COLOR,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size_pt)
        run.font.color.rgb = color
        if bold_first and i == 0:
            run.font.bold = True
    return box


def add_two_column_slide(
    slide,
    left_title: str,
    left_lines: list[str],
    right_title: str,
    right_lines: list[str],
) -> None:
    margin_l = Inches(0.55)
    col_w = Inches(5.6)
    top = Inches(2.0)
    height = Inches(4.8)
    gap = Inches(0.35)
    right_l = margin_l + col_w + gap

    add_textbox(
        slide,
        margin_l,
        top,
        col_w,
        Inches(0.45),
        [left_title],
        size_pt=24,
        bold_first=True,
        color=ACCENT_QPU,
    )
    add_textbox(slide, margin_l, top + Inches(0.5), col_w, height, left_lines, size_pt=20)

    add_textbox(
        slide,
        right_l,
        top,
        col_w,
        Inches(0.45),
        [right_title],
        size_pt=24,
        bold_first=True,
        color=ACCENT_SIM,
    )
    add_textbox(slide, right_l, top + Inches(0.5), col_w, height, right_lines, size_pt=20)


def remove_shapes_by_name(slide, names: set[str]) -> None:
    to_remove = [shape for shape in slide.shapes if shape.name in names]
    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)


def clear_protocol_slide(slide) -> None:
    remove_shapes_by_name(slide, {"TextBox 20", "Picture 21"})


def update_slide_16_bridge(slide) -> None:
    clear_protocol_slide(slide)
    remove_shapes_by_name(slide, {"TextBox 19"})

    set_title(find_shape(slide, "TextBox 8"), "Expressibility — two complementary views")
    set_footer(find_shape(slide, "TextBox 7"), "10 / 16")

    add_textbox(
        slide,
        Inches(0.55),
        Inches(2.05),
        Inches(11.8),
        Inches(0.55),
        ["How broadly does the ansatz explore state space?"],
        size_pt=24,
        color=MUTED_COLOR,
    )

    add_two_column_slide(
        slide,
        "QPU track",
        [
            "Expressibility on real IQM Spark hardware",
            "Includes noise, drift, and tomography error",
            "Answers: what actually happens on device?",
            "Metric: D_KL(P_QPU || P_Haar)",
        ],
        "Simulator track",
        [
            "Ideal noiseless expressibility (statevector)",
            "Upper bound for the same circuit definitions",
            "Answers: what could the ansatz do without noise?",
            "Metric: D_KL(P_Sim || P_Haar)",
        ],
    )

    add_textbox(
        slide,
        Inches(0.55),
        Inches(6.55),
        Inches(11.8),
        Inches(0.9),
        [
            "Same random parameter pairs, ansätze, depths, bins, and KL definition "
            "→ results are directly comparable."
        ],
        size_pt=20,
        color=MUTED_COLOR,
    )


def update_slide_17_shared(slide) -> None:
    clear_protocol_slide(slide)

    set_title(find_shape(slide, "TextBox 8"), "KL expressibility — shared protocol")
    set_footer(find_shape(slide, "TextBox 7"), "10 / 16")

    ref = find_shape(slide, "TextBox 19")
    if ref is not None:
        set_single_line_text(
            ref,
            "Sim et al., Adv. Quantum Technol. 2019  ·  "
            "Haar reference: Olivia Di Matteo",
            size_pt=14,
        )

    add_textbox(
        slide,
        Inches(0.55),
        Inches(2.15),
        Inches(11.8),
        Inches(3.2),
        [
            "Shared steps (both QPU and Simulator):",
            "1. Sample random parameter pairs (θ_a, θ_b) — matched seeds",
            "2. For each ansatz and depth, collect pairwise state overlaps F",
            "3. Build empirical histogram P(F) and compare to analytic P_Haar(F)",
            "4. Expressibility score: D_KL(P_emp || P_Haar)  —  lower is better",
        ],
        size_pt=24,
    )


def update_slide_18_dual_paths(slide) -> None:
    clear_protocol_slide(slide)

    set_title(find_shape(slide, "TextBox 8"), "KL expressibility — QPU vs Simulator")
    set_footer(find_shape(slide, "TextBox 7"), "10 / 16")

    ref = find_shape(slide, "TextBox 19")
    if ref is not None:
        set_single_line_text(
            ref,
            "Same ansätze (ansatz_odra, ansatz_simulator) and depths on both tracks.",
            size_pt=14,
        )

    add_two_column_slide(
        slide,
        "QPU (IQM Spark)",
        [
            "Prepare two states on hardware",
            "Full Pauli tomography (3^n bases)",
            "F = Tr(ρ_a ρ_b) with shot noise",
            "Costly — pilot-chosen shots, pairs, iterations",
            "D_KL(P_QPU || P_Haar)",
        ],
        "Simulator (statevector)",
        [
            "Evolve two ideal statevectors",
            "Exact overlap |⟨ψ_a|ψ_b⟩|²",
            "Noiseless fidelities",
            "Cheap — thousands of pairs (e.g. 5000)",
            "D_KL(P_Sim || P_Haar)",
        ],
    )


def update_slide_19_convergence(slide) -> None:
    clear_protocol_slide(slide)

    set_title(find_shape(slide, "TextBox 8"), "KL expressibility — from fidelities to score")
    set_footer(find_shape(slide, "TextBox 7"), "10 / 16")

    ref = find_shape(slide, "TextBox 19")
    if ref is not None:
        set_single_line_text(ref, "Both tracks converge to the same analysis pipeline.", size_pt=14)

    add_textbox(
        slide,
        Inches(0.55),
        Inches(2.0),
        Inches(11.8),
        Inches(1.2),
        [
            "θ pairs  →  QPU: tomography on Spark  |  Sim: statevector evolution",
            "           →  pairwise fidelities F  →  histogram P(F)",
            "           →  compare to P_Haar  →  KL score",
        ],
        size_pt=22,
        color=MUTED_COLOR,
    )

    add_textbox(
        slide,
        Inches(0.55),
        Inches(3.5),
        Inches(11.8),
        Inches(2.5),
        [
            "• Same n_bins, ε-smoothing, and bootstrap confidence intervals",
            "• Both ansätze tested at the same depths",
            "• Sim: both ansätze approach Haar at depth 4; QPU shows a larger odra–simulator gap",
        ],
        size_pt=24,
    )


def main() -> None:
    shutil.copy2(PPTX, BACKUP)
    prs = Presentation(str(PPTX))

    update_slide_16_bridge(prs.slides[15])
    update_slide_17_shared(prs.slides[16])
    update_slide_18_dual_paths(prs.slides[17])
    update_slide_19_convergence(prs.slides[18])

    prs.save(str(PPTX))
    print(f"Saved: {PPTX}")
    print(f"Backup: {BACKUP}")
    print(f"Slide count: {len(prs.slides)}")


if __name__ == "__main__":
    main()
